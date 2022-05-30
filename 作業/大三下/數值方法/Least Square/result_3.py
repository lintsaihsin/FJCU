import matplotlib.pyplot as plt
import argparse
import numpy as np
import collections.abc
from pptx import Presentation, util
from pptx.util import Pt
import os
import pandas as pd

x = []
y = []
sigx = []
sigxy = []
mse = []

def sig1(t): # sigma(Xi 的 t次方)
    summ = 0
    for i in range(n):
        summ += x[i]**t
    return summ

def sig2(t): # sigma(Yi * (Xi 的 t次方))
    summ = 0
    for i in range(n):
        summ += y[i] * x[i]**t
    return summ

try:
    while(1):
        inp = input().split(" ")
        x.append(float(inp[0]))
        y.append(float(inp[1]))
except EOFError:
    pass
n = len(x)

for i in range(n):
    print(x[i])
for i in range(n):
    print(y[i])

sigx.append(n) # n : 實驗個數
for i in range(1, n*2+1):
    sigx.append(sig1(i)) # sigma(Xi^t)
    sigxy.append(sig2(i-1)) # sigma(Yi * (Xi^t))
    print(format(sig1(i), '.4E'), " ", format(sig2(i-1), '.4E'))

name = np.array([ "" for _ in range(i) ], dtype=object) # 完整的方程式

for i in range(1, n):
    print("Case ", i, " : ")
    matrix = np.array([ [ 0. for _ in range(i) ] for _ in range(i) ])
    matrixy = []
    for j in range(i):
        matrixy.append(sigxy[j])
        for k in range(i):
            matrix[j][k] = sigx[j+k]
    #=========matrix=========
    # n         sig(Xi)   sig(Xi^2) ...
    # sig(Xi)   sig(Xi^2) sig(Xi^3) ...
    # sig(Xi^2) sig(Xi^3) sig(Xi^4) ...
    # ...       ...       ...
    #========================
    #=========matrixy=========
    # sig(Yi) 
    # sig(YiXi) 
    # sig(Yi * (Xi^2)) 
    # ...
    #========================
    matrix = np.array(matrix)
    matrixy = np.array(matrixy)
    ans = np.linalg.solve(matrix, matrixy)
    #=========ans=========
    # a0 
    # a1 
    # a2 
    # ...
    #=====================
    print(ans)
    pm = np.array([ 0. for _ in range(n) ])
    for j in range(n):
        sum_y = 0
        m = 1
        for k in range(i):
            sum_y += m * ans[k]
            m *= x[j]
        pm[j] = sum_y
    mse.append((np.sum((pm-y) ** 2)/(n-i))**0.5) # Mean square error

#==================================畫線========================================================
    _x = np.linspace(min(x), max(x), 10000) #放入10000個範圍在 -10 ~ 10的x座標
    _y = np.array([ 0. for _ in range(10000) ])
    for j in range(10000):
        sum_y = 0
        m = 1
        for k in range(i):
            sum_y += m * ans[k]
            m *= _x[j]
        _y[j] = sum_y
#==================================完整的方程式=================================================        
    name[i] = "f(x) = "
    for k in range(i):
        if k > 1:
            name[i] += " + " + format(ans[k], '.4E') + " * x^" + str(k)
        elif k == 1:
            name[i] += " + " + format(ans[k], '.4E') + " * x"
        else :
            name[i] += format(ans[k], '.4E')
#==================================畫圖========================================================
    print("") #title
    print("=" * 50)
    plt.ylabel("y") # y label
    plt.xlabel("x") # x label
    plt.plot(_x, _y)
    for j in range(n):
        p1 = plt.scatter(x[j-1], y[j-1] ,color='black', marker = 'o')
    # plt.legend([p1, p2, p3, p4], ['+1(Training)', '- 1(Training)', '+1(Testing)', '- 1(Testing)'], loc = 'lower right')
    plt.grid() #格線
    if i < 10:
        plt.savefig(f'./Result_3/Case_0{i}.png')
    else :
        plt.savefig(f'./Result_3/Case_{i}.png')
    plt.clf()
#=============================================================================================    
find_min = np.array([ 0. for _ in range(n-2) ])
print("mse size : ", len(mse))
print("mse")
for i in range(len(mse)):
    print(mse[i])
    
bc = [] # best choice
print("find min")
for i in range(n-2):
    if mse[i] / mse[i+1] < 2:
        bc.append(i+1)
    find_min[i] = mse[i] / mse[i+1] #sig(k) / sig(k+1)
    print(find_min[i])

print(find_min)
print("Case 1 :", np.argmin(find_min)+1)
print("sort :", np.argsort(find_min)+1)
print("THE BEST CHOICE : ", bc)
print("Case 2 :", np.argmin(mse)+1)
print("sort :", np.argsort(mse)+1)
print("name", name)


#==================================ppt========================================================
def add_slide(prs, pic_list, result_path, file, path , number):
    # 範例化空白模板
    blank_slide_layout = prs.slide_layouts[6]
    # 向檔案中新增空白頁面
    slide = prs.slides.add_slide(blank_slide_layout)

    width = util.Cm(5)
    height = util.Cm(5)
    for p in pic_list:
        # 圖片路徑
        img_path = f'{result_path}/' + p
        # 設定圖片位置
        n = pic_list.index(p)
        top = util.Cm((n//2) * 7.15 + 0.35)
        left = util.Cm((n%2) * 13)
        pic = slide.shapes.add_picture(img_path, left, top, width, height)
    
    left = top = width = height = util.Cm(1)
    print("range : ", number-4, " ",  number)
    for i in range(number-4, number):
        top = util.Cm(((i+4-number)//2) * 7.15)
        left = util.Cm(((i+4-number)%2) * 13)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "Case " + str(i+1) + " :"
        # tf.font.size = Pt(18)

        # if i != 43:
        top = util.Cm(((i+4-number)//2) * 7 + 5.25)
        left = util.Cm(((i+4-number)%2) * 13 + 0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        print("error : ", i)
        tf.text = format(mse[i], '.4E')
        # tf.font.size = Pt(10)

        if i != 43:
            top = util.Cm(((i+4-number)//2) * 7 + 6)
            left = util.Cm(((i+4-number)%2) * 13 + 0.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = format(find_min[i], '.4E')
            # tf.font.size = Pt(10)

        width = util.Cm(7.5)
        height = util.Cm(6.5)
        top = util.Cm(((i+4-number)//2) * 7 + 0.35)
        left = util.Cm(((i+4-number)%2) * 13 + 5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = name[i+1]
        # tf.font.size = Pt(8)

    return slide


# 範例化一個ppt簡報物件
prs = Presentation()
prs.slide_width = util.Cm(25.4)
prs.slide_height = util.Cm(14.29)
result_path = "./Result_3/"
# 讀取圖片列表
pic_list = []
tmp = 0
path = f'{result_path}'
for file in sorted(os.listdir(result_path)):
    if '.png' in file:
        pic_list.append(file)
        tmp += 1  
    print(file)
    if len(pic_list) % 4 == 0:
        slide = add_slide(prs, pic_list, result_path, file, path , tmp)
        pic_list = []

# slide = add_slide(prs, pic_list, result_path, file, path, name, len(pic_list))
print('圖片列表 : ', pic_list, len(pic_list))
print("="*50)


        

# 設定圖片的大小


# 儲存為檔案
prs.save("Result_3.pptx")